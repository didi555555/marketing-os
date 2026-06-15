from flask import Blueprint, render_template, request, send_file, current_app
from flask_login import login_required
import io, json, random, re, math
from datetime import datetime
from PIL import Image, ImageDraw, ImageFont, ImageFilter
try:
    import cairosvg
    HAS_CAIRO = True
except:
    HAS_CAIRO = False
HAS_SVG_RENDER = False
try:
    import xml.etree.ElementTree as ET_xml
    from svg.path import parse_path as svg_parse_path
    from svg.path.path import CubicBezier as CB, QuadraticBezier as QB, Arc as SvgArc
    import re as svg_re
    HAS_SVG_RENDER = True
except:
    pass

bp = Blueprint('media_studio', __name__, template_folder='templates',
               url_prefix='/plugin/media-studio')

@bp.route('/')
@login_required
def index():
    return render_template('media_index.html')

@bp.route('/generate-image', methods=['POST'])
@login_required
def generate_image():
    prompt = request.form.get('prompt', 'تسويق')
    width, height = 800, 600
    img_bytes = try_gemini_svg(prompt, width, height)
    if img_bytes:
        return send_file(io.BytesIO(img_bytes), mimetype='image/png', as_attachment=False)
    buf = render_fallback_image(prompt, width, height)
    return send_file(buf, mimetype='image/png', as_attachment=False)

def parse_svg_style(elem, key, default=''):
    style = elem.get('style', '')
    if style:
        m = re.search(r'(?:^|;)\s*' + key + r'\s*:\s*([^;]+)', style)
        if m:
            return m.group(1).strip()
    return elem.get(key, default)

def svg_color_to_rgb(c, default=(128,128,128)):
    c = c.strip()
    if c.startswith('url(#'):
        return None, c[5:-1]
    if c.startswith('#'):
        c = c[1:]
        if len(c) == 3: c = ''.join(t*2 for t in c)
        return tuple(int(c[i:i+2],16) for i in (0,2,4)), None
    if c == 'none': return None, None
    nums = re.findall(r'\d+', c)
    if len(nums) >= 3: return tuple(int(n) for n in nums[:3]), None
    return default, None

def render_svg_to_png(svg_text, out_width=800, out_height=600):
    if not HAS_SVG_RENDER:
        return None
    try:
        import io as svg_io
        root = ET_xml.fromstring(svg_text)
        grads = {}
        for d in root.iter():
            if d.tag.endswith('}defs'):
                for g in d:
                    gid = g.get('id')
                    if gid and ('Gradient' in g.tag):
                        stops = []
                        for s in g:
                            off = float(s.get('offset','0').replace('%',''))
                            c_str = parse_svg_style(s, 'stop-color', '#888')
                            c, _ = svg_color_to_rgb(c_str)
                            if c: stops.append((off, c))
                        if stops: grads[gid] = stops
        
        def resolve_color(c_str):
            c, gid = svg_color_to_rgb(c_str)
            if c: return c
            if gid in grads and grads[gid]:
                colors = [s[1] for s in grads[gid]]
                return tuple(sum(clr[i] for clr in colors)//len(colors) for i in range(3))
            return (200, 200, 150)
        
        def get_svg_points(d):
            pts = []
            p = svg_parse_path(d)
            for seg in p:
                n = 24 if isinstance(seg, (CB, QB, SvgArc)) else 2
                for i in range(n+1):
                    pt = seg.point(i/n)
                    pts.append((pt.real, pt.imag))
            return pts
        
        def transform_pts(pts, dx=0, dy=0, angle=0):
            if not pts: return pts
            cos_a, sin_a = math.cos(angle), math.sin(angle)
            return [(x*cos_a - y*sin_a + dx, x*sin_a + y*cos_a + dy) for x,y in pts]
        
        img = Image.new('RGBA', (out_width, out_height), (0,0,0,0))
        draw = ImageDraw.Draw(img)
        
        def draw_svg_elem(e, dx=0, dy=0, angle=0):
            tag = e.tag.split('}')[-1]
            t = e.get('transform','')
            if t:
                m = re.search(r'translate\(([^)]+)\)', t)
                if m:
                    parts = m.group(1).replace(',',' ').split()
                    dx += float(parts[0])
                    if len(parts) > 1: dy += float(parts[1])
                m = re.search(r'rotate\(([^)]+)\)', t)
                if m: angle += float(m.group(1)) * math.pi/180
            if tag == 'g':
                for c in e: draw_svg_elem(c, dx, dy, angle)
                return
            fill = resolve_color(parse_svg_style(e, 'fill', ''))
            stroke = resolve_color(parse_svg_style(e, 'stroke', ''))
            sw = float(parse_svg_style(e, 'stroke-width', '1'))
            if tag == 'rect':
                x = float(e.get('x',0)) + dx
                y = float(e.get('y',0)) + dy
                rw = float(e.get('width',0))
                rh = float(e.get('height',0))
                rx = float(e.get('rx',0))
                if fill:
                    if rx: draw.rounded_rectangle([x,y,x+rw,y+rh], radius=rx, fill=fill)
                    else: draw.rectangle([x,y,x+rw,y+rh], fill=fill)
                if stroke: draw.rectangle([x,y,x+rw,y+rh], outline=stroke, width=max(1,int(sw)))
            elif tag == 'path':
                d = e.get('d','')
                pts = transform_pts(get_svg_points(d), dx, dy, angle)
                if stroke and len(pts) > 1:
                    wd = max(1, int(sw))
                    for i in range(len(pts)-1):
                        draw.line([pts[i], pts[i+1]], fill=stroke, width=wd)
                if fill and len(pts) > 2:
                    draw.polygon(pts, fill=fill)
            elif tag == 'ellipse':
                cx = float(e.get('cx',0)) + dx
                cy = float(e.get('cy',0)) + dy
                rx2 = float(e.get('rx',0))
                ry2 = float(e.get('ry',0))
                if fill: draw.ellipse([cx-rx2,cy-ry2,cx+rx2,cy+ry2], fill=fill)
            elif tag == 'line':
                x1 = float(e.get('x1',0)) + dx
                y1 = float(e.get('y1',0)) + dy
                x2 = float(e.get('x2',0)) + dx
                y2 = float(e.get('y2',0)) + dy
                if stroke: draw.line([x1,y1,x2,y2], fill=stroke, width=max(1,int(sw)))
            elif tag == 'circle':
                cx = float(e.get('cx',0)) + dx
                cy = float(e.get('cy',0)) + dy
                r = float(e.get('r',0))
                if fill: draw.ellipse([cx-r,cy-r,cx+r,cy+r], fill=fill)
        
        for child in root:
            if not child.tag.endswith('}defs'):
                draw_svg_elem(child)
        
        bg = Image.new('RGB', (out_width, out_height), (100, 80, 60))
        bg.paste(img, (0,0), img)
        buf = io.BytesIO()
        bg.save(buf, 'PNG')
        buf.seek(0)
        return buf.read()
    except:
        return None

def try_gemini_svg(prompt, width, height):
    api_key = current_app.config.get('GEMINI_API_KEY', '')
    if not api_key:
        return None
    try:
        import requests
        url = 'https://generativelanguage.googleapis.com/v1/models/gemini-2.5-flash-lite:generateContent'
        headers = {'X-Goog-Api-Key': api_key, 'Content-Type': 'application/json'}
        data = {
            'contents': [{'parts': [{'text': f'Generate ONLY valid SVG code (no markdown, no explanation, no wrapper) for an image: {prompt}. Size {width}x{height}. Use gradients, shadows, realistic colors. Return ONLY the raw SVG starting with <svg and ending with </svg>'}]}],
            'generationConfig': {'temperature': 0.4, 'maxOutputTokens': 4000}
        }
        r = requests.post(url, json=data, headers=headers, timeout=25)
        if r.status_code != 200:
            return None
        text = r.json()['candidates'][0]['content']['parts'][0]['text']
        match = re.search(r'<svg[\s\S]*?</svg>', text)
        if not match:
            return None
        svg = match.group(0)
        return render_svg_to_png(svg, width, height)
    except:
        return None

def render_fallback_image(prompt, width, height):
    img = Image.new('RGB', (width, height), (60, 50, 40))
    draw = ImageDraw.Draw(img)
    for y in range(height):
        r = int(100 - y * 0.04)
        g = int(70 - y * 0.03)
        b = int(35 - y * 0.02)
        draw.line([(0, y), (width, y)], fill=(max(30,r), max(20,g), max(10,b)))
    
    is_cheese = any(w in prompt.lower() for w in ['جبن', 'جبنة', 'cheese', 'جبنه'])
    if is_cheese:
        colors = [(210, 170, 50), (195, 155, 40), (225, 190, 60)]
        for layer in range(3):
            offset_x = 150 + layer * 30
            offset_y = 180 + layer * 25
            for i in range(7 - layer):
                x = offset_x + i * 75
                y = offset_y
                c = colors[layer % len(colors)]
                shade = random.randint(-20, 20)
                c = (c[0]+shade, c[1]+shade, c[2]-10)
                draw.rounded_rectangle([x, y, x+55, y+180-layer*20], radius=14, fill=c)
                draw.rounded_rectangle([x+5, y+5, x+50, y+175-layer*20], radius=11, fill=(min(255,c[0]+30), min(255,c[1]+35), c[2]+5))
                for j in range(6):
                    ly = y + 20 + j * 28
                    draw.line([(x+8, ly), (x+47, ly+15)], fill=(c[0]-40, c[1]-30, 10), width=3)
        shadow = Image.new('L', (width, height), 0)
        sd = ImageDraw.Draw(shadow)
        sd.ellipse([120, 480, 680, 530], fill=100)
        shadow = shadow.filter(ImageFilter.GaussianBlur(10))
        img = Image.composite(img, Image.new('RGB', (width, height), (0,0,0)), shadow)
        draw = ImageDraw.Draw(img)
    else:
        for i in range(30):
            x1 = random.randint(0, width)
            y1 = random.randint(0, height)
            x2 = x1 + random.randint(-80, 80)
            y2 = y1 + random.randint(-80, 80)
            x1, x2 = min(x1, x2), max(x1, x2)
            y1, y2 = min(y1, y2), max(y1, y2)
            draw.arc([x1, y1, x2, y2], 0, 360, fill=(random.randint(100,255), random.randint(100,255), random.randint(100,255)))
    
    draw.text((width//2, 30), prompt[:50], fill=(255, 240, 200), anchor='mt')
    draw.text((width//2, height-15), 'Marketing OS', fill=(180, 160, 120), anchor='mb')
    buf = io.BytesIO()
    img.save(buf, 'PNG')
    buf.seek(0)
    return buf

@bp.route('/generate-docx', methods=['POST'])
@login_required
def generate_docx():
    from docx import Document
    title = request.form.get('title', 'تقرير تسويقي')
    body = request.form.get('body', 'محتوى التقرير')
    doc = Document()
    doc.add_heading(title, 0)
    doc.add_paragraph(body)
    doc.add_paragraph(f'تم الإنشاء: {datetime.now()}')
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return send_file(buf, as_attachment=True, download_name='report.docx',
                     mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document')

@bp.route('/generate-pdf', methods=['POST'])
@login_required
def generate_pdf():
    from fpdf import FPDF
    title = request.form.get('title', 'تقرير')
    body = request.form.get('body', 'محتوى التقرير')
    pdf = FPDF()
    pdf.add_page()
    pdf.add_font('Arabic', '', 'C:\\Windows\\Fonts\\arabtype.ttf')
    pdf.set_font('Arabic', '', 16)
    pdf.cell(0, 10, title, new_x='LMARGIN', new_y='NEXT', align='C')
    pdf.set_font('Arabic', '', 12)
    pdf.multi_cell(0, 10, body)
    buf = io.BytesIO()
    pdf.output(buf)
    buf.seek(0)
    return send_file(buf, as_attachment=True, download_name='report.pdf', mimetype='application/pdf')

@bp.route('/generate-pptx', methods=['POST'])
@login_required
def generate_pptx():
    from pptx import Presentation
    from pptx.util import Inches
    title = request.form.get('title', 'عرض تقديمي')
    slides_text = request.form.get('slides', 'شريحة 1\nنقطة 1\nنقطة 2')
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    lines = slides_text.split('\n')
    content = '\n'.join(lines[1:]) if len(lines) > 1 else slides_text
    slide.placeholders[1].text = content if content else title
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return send_file(buf, as_attachment=True, download_name='presentation.pptx',
                     mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

@bp.route('/generate-xlsx', methods=['POST'])
@login_required
def generate_xlsx():
    from openpyxl import Workbook
    data = request.form.get('data', 'المنتج,السعر,الكمية\nمنتج أ,100,5\nمنتج ب,200,3')
    wb = Workbook()
    ws = wb.active
    ws.title = 'بيانات'
    for row_data in data.strip().split('\n'):
        ws.append(row_data.split(','))
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return send_file(buf, as_attachment=True, download_name='data.xlsx',
                     mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')

def register(app):
    app.register_blueprint(bp)
